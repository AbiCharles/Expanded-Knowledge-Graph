"""Generate docs/overview.pptx — a 4-slide manager-facing summary deck.

Slide 1 — Audiences and surfaces  (Operator / Reviewer / Engineer, each with
          what they touch, what they get, and an example workflow snippet)
Slide 2 — Architecture: diagram + mental model  (existing diagram on the left,
          the three nouns + three verbs stacked on the right)
Slide 3 — How scenarios · ontology · actions · KnowledgeResolver interact
          (mini-diagram of the four pieces, captions for each, fallback panel)
Slide 4 — How a request flows through the system  (numbered walk-through of
          one operator prompt, tagging which layer activates at each step)

Run with:
    .venv/bin/python scripts/build_overview_pptx.py
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.patches as patches
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
DIAGRAM_PNG = REPO / "docs" / "_architecture-diagram.png"
COMPONENT_PNG = REPO / "docs" / "_component-interaction.png"
OUT_PPTX = REPO / "docs" / "overview.pptx"

# Colors
NAVY = RGBColor(0x14, 0x26, 0x47)
NAVY_FILL = RGBColor(0xE7, 0xED, 0xF6)
CYAN_FILL = RGBColor(0xDF, 0xF7, 0xF3)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD1)
EMERALD_FILL = RGBColor(0xDF, 0xF1, 0xE7)
GREY = RGBColor(0x55, 0x66, 0x77)
LIGHT_GREY = RGBColor(0xF3, 0xF5, 0xF9)
INK = RGBColor(0x23, 0x35, 0x55)


def _add_text_box(slide, left, top, width, height, *, fill=None, border=False):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if border:
        box.line.color.rgb = NAVY
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    return tf


def _set_run(run, text, *, size=11, bold=False, color=INK, italic=False,
             name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def _title(slide, text, *, top=Inches(0.3)):
    tf = _add_text_box(slide, Inches(0.4), top, Inches(12.6), Inches(0.6))
    p = tf.paragraphs[0]
    _set_run(p.add_run(), text, size=24, bold=True, color=NAVY)


def _subtitle(slide, text, *, top=Inches(0.95)):
    tf = _add_text_box(slide, Inches(0.4), top, Inches(12.6), Inches(0.4))
    p = tf.paragraphs[0]
    _set_run(p.add_run(), text, size=12, italic=True, color=GREY)


def _bullet(tf, text, *, level=0, size=11, color=INK, bold_lead=False):
    p = tf.add_paragraph()
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    if bold_lead and " — " in text:
        lead, rest = text.split(" — ", 1)
        r1 = p.add_run()
        _set_run(r1, "• " + lead, size=size, bold=True, color=NAVY)
        r2 = p.add_run()
        _set_run(r2, " — " + rest, size=size, color=color)
    else:
        r = p.add_run()
        _set_run(r, "• " + text, size=size, color=color)
    p.space_after = Pt(3)


# =============================================================================
# Slide 1 — Audiences and surfaces
# =============================================================================
def slide_audiences(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide, "Three audiences, three surfaces")
    _subtitle(slide,
              "Each persona has its own surface; nobody is asked to learn "
              "another role's tools. Examples of real workflow snippets below.")

    cards = [
        {
            "title": "Operator",
            "touches": "Chat composer + a row of suggested-prompt chips",
            "gets": "A natural-language way to kick off pre-approved workflows or "
                    "ad-hoc data queries against any registered source. No SQL or "
                    "ontology authoring required.",
            "example_label": "Example prompt",
            "example": "“Override the SC-TC-001 sanctions block on order "
                       "ORD-44216 — customer says they have an OFAC license.”",
            "fill": AMBER_FILL,
        },
        {
            "title": "Reviewer",
            "touches": "A Teams Adaptive Card (or in-app card) per pending decision",
            "gets": "Approve / Reject / Need more info with the full bound "
                    "context: active policy, actor scope, master data, prior "
                    "similar cases, SOPs — all attached automatically.",
            "example_label": "Example card",
            "example": "Sanctioned counterparty: Sanctioned Pharma Holdings · "
                       "Product P-EL-9001 (ECCN 5A002) · Contract €4.2M "
                       "→ [Approve] [Reject] [Need info]",
            "fill": CYAN_FILL,
        },
        {
            "title": "Engineer / scenario author",
            "touches": "YAML files + the Knowledge admin tile",
            "gets": "Register a new system and every existing scenario picks it "
                    "up through the ontology layer — no rewrites. Author one "
                    "ontology, get many scenarios for free.",
            "example_label": "Example change",
            "example": "Add a Postgres governance source → click “Suggest "
                       "mappings” → confirm Supplier.country → "
                       "suppliers.country_code → done.",
            "fill": EMERALD_FILL,
        },
    ]

    left_positions = [Inches(0.4), Inches(4.7), Inches(9.0)]
    card_w = Inches(4.0)
    card_h = Inches(5.7)

    for left, card in zip(left_positions, cards):
        tf = _add_text_box(slide, left, Inches(1.55), card_w, card_h,
                           fill=card["fill"], border=True)

        # Card title
        p = tf.paragraphs[0]
        _set_run(p.add_run(), card["title"], size=18, bold=True, color=NAVY)
        p.space_after = Pt(8)

        # "What they touch"
        p = tf.add_paragraph()
        _set_run(p.add_run(), "What they touch", size=10, bold=True, color=GREY)
        p.space_after = Pt(2)
        p = tf.add_paragraph()
        _set_run(p.add_run(), card["touches"], size=11.5, color=INK)
        p.space_after = Pt(8)

        # "What they get"
        p = tf.add_paragraph()
        _set_run(p.add_run(), "What they get", size=10, bold=True, color=GREY)
        p.space_after = Pt(2)
        p = tf.add_paragraph()
        _set_run(p.add_run(), card["gets"], size=11.5, color=INK)
        p.space_after = Pt(10)

        # Example snippet
        p = tf.add_paragraph()
        _set_run(p.add_run(), card["example_label"], size=10, bold=True, color=GREY)
        p.space_after = Pt(2)
        p = tf.add_paragraph()
        _set_run(p.add_run(), card["example"], size=11, italic=True, color=NAVY,
                 name="Consolas")

    # Footer line
    tf = _add_text_box(slide, Inches(0.4), Inches(7.0), Inches(12.6), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(),
             "Same case object travels across all three surfaces — the "
             "audit trail stitches them together.",
             size=11, italic=True, color=GREY)


# =============================================================================
# Slide 2 — Architecture diagram + mental model
# =============================================================================
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide, "Architecture at a glance")
    _subtitle(slide,
              "Top of the diagram: what the operator sees. Bottom: connected "
              "enterprise systems. Each row is a swap-point.")

    # Diagram on the left
    slide.shapes.add_picture(
        str(DIAGRAM_PNG),
        Inches(0.4), Inches(1.45),
        width=Inches(8.4),
    )

    # Right-hand "Mental model" panel
    tf = _add_text_box(slide, Inches(9.0), Inches(1.45), Inches(4.0),
                       Inches(5.85), fill=LIGHT_GREY, border=True)
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "The mental model", size=16, bold=True, color=NAVY)
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    _set_run(p.add_run(), "Three nouns + three verbs.",
             size=10.5, italic=True, color=GREY)
    p.space_after = Pt(8)

    # Nouns
    p = tf.add_paragraph()
    _set_run(p.add_run(), "Three nouns  (what it knows)",
             size=11.5, bold=True, color=NAVY)
    p.space_after = Pt(2)
    _bullet(tf, "Data sources — CSV, SQLite, Postgres, HTTP, vector, Neo4j",
            size=10.5, bold_lead=True)
    _bullet(tf, "Ontology — business classes mapped to source columns; authored once, reused everywhere",
            size=10.5, bold_lead=True)
    _bullet(tf, "Scenarios — YAML recipes the agent follows; one per pattern",
            size=10.5, bold_lead=True)

    # Verbs
    p = tf.add_paragraph()
    _set_run(p.add_run(), "Three verbs  (what happens)",
             size=11.5, bold=True, color=NAVY)
    p.space_before = Pt(10)
    p.space_after = Pt(2)
    _bullet(tf, "Classify — LLM picks the scenario that matches the prompt",
            size=10.5, bold_lead=True)
    _bullet(tf, "Bind — framework gathers facts with provenance into a typed envelope",
            size=10.5, bold_lead=True)
    _bullet(tf, "Decide — auto-execute (low risk) or pause for a named reviewer",
            size=10.5, bold_lead=True)

    # Closing line
    p = tf.add_paragraph()
    p.space_before = Pt(12)
    _set_run(p.add_run(),
             "Everything else is implementation.",
             size=10.5, italic=True, color=GREY)


# =============================================================================
# Component-interaction mini-diagram  (rendered by matplotlib, saved to PNG)
# =============================================================================
def render_component_diagram(out: Path) -> None:
    """Render the 4-way interaction diagram as a clean two-column flow.

    Left column  (read path):   Scenario → Ontology → KnowledgeResolver → bound facts
    Right column (write path):  Scenario → Action  → Executors
    """
    fig, ax = plt.subplots(figsize=(12, 7), dpi=160)
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 11)
    ax.axis("off")

    NAVY = "#142647"
    NAVY_FILL = "#e7edf6"
    CYAN_FILL = "#dff7f3"
    AMBER_FILL = "#fbeed1"
    EMERALD_FILL = "#dff1e7"
    WHITE = "#ffffff"

    def box(x, y, w, h, label, *, fill, edge=NAVY, fontsize=10, weight="bold"):
        rect = patches.FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.02,rounding_size=0.08",
            linewidth=1.4, edgecolor=edge, facecolor=fill,
        )
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, label,
                ha="center", va="center",
                fontsize=fontsize, color=NAVY, fontweight=weight)

    def vert_arrow(x, y_top, y_bot, *, label=None, color="#395270", lw=1.7,
                   label_x_offset=0.0, label_size=9.5, style="-|>"):
        # Arrow points downward
        ax.annotate(
            "", xy=(x, y_bot), xytext=(x, y_top),
            arrowprops=dict(arrowstyle=style, color=color, lw=lw),
        )
        if label:
            ax.text(x + label_x_offset, (y_top + y_bot) / 2, label,
                    fontsize=label_size, color="#23344a",
                    ha="left" if label_x_offset >= 0 else "right",
                    va="center", style="italic")

    # === Top: Scenario (centre, spans both columns) ===
    box(3.0, 9.0, 10.0, 1.6,
        "Scenario  (YAML recipe)\n"
        "stages: facts + ontology_queries  ·  who reviews  ·  what to execute on approve",
        fill=CYAN_FILL, fontsize=11)

    # === Middle row: Ontology (left)   |   Action (right) ===
    box(0.6, 5.8, 6.4, 2.4,
        "Ontology + mappings\n\n"
        "Business classes — Product,\n"
        "SanctionedEntity, Supplier — mapped\n"
        "to columns in registered sources.\n"
        "Authored once, reused everywhere.",
        fill=CYAN_FILL, fontsize=10, weight="normal")

    box(9.0, 5.8, 6.4, 2.4,
        "Action  (typed write operation)\n\n"
        "name · typed arg schema ·\n"
        "named executor reference.\n"
        "HITL by default — LLM only fills args;\n"
        "executor runs after reviewer approves.",
        fill=AMBER_FILL, fontsize=10, weight="normal")

    # === Bottom row: KnowledgeResolver (left) | Executors (right) ===
    box(0.6, 2.7, 6.4, 2.0,
        "KnowledgeResolver Protocol\n\n"
        "CSV · SQLite · Postgres · HTTP ·\n"
        "Vector store · Neo4j   (all swappable)",
        fill=NAVY_FILL, fontsize=10)

    box(9.0, 2.7, 6.4, 2.0,
        "Executors\n\n"
        "sql_update · http_request\n"
        "(named, hand-written; LLM never picks them)",
        fill=EMERALD_FILL, fontsize=10)

    # === Bottom-left: Bound facts envelope (output of the read path) ===
    box(2.6, 0.2, 6.4, 1.8,
        "Bound facts  (typed envelope)\n\n"
        "every fact carries provenance:\n"
        "which ontology class asked · which source answered",
        fill=WHITE, fontsize=9.5, weight="normal")

    # === Vertical arrows — left column (read path) ===
    vert_arrow(3.8, 8.95, 8.25,
               label="  references ontology class",
               label_x_offset=0.2)
    vert_arrow(3.8, 5.75, 4.75,
               label="  resolve_query(class, where)",
               label_x_offset=0.2)
    vert_arrow(3.8, 2.65, 2.05,
               label="  rows → facts",
               label_x_offset=0.2, color="#3b6e58")

    # === Vertical arrows — right column (write path) ===
    vert_arrow(12.2, 8.95, 8.25,
               label="  on approve → execute stage",
               label_x_offset=0.2)
    vert_arrow(12.2, 5.75, 4.75,
               label="  executor.run(args)",
               label_x_offset=0.2)

    # Caption strip at bottom
    ax.text(0.6, -0.3,
            "Read path (left): scenarios reference ontology classes; "
            "the resolver translates those references into rows.   "
            "Write path (right): actions are how the system writes back — only after a named reviewer approves.",
            fontsize=9, color="#444", style="italic")

    plt.tight_layout()
    plt.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# =============================================================================
# Slide 3 — Component interaction  (NEW)
# =============================================================================
def slide_components(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide,
           "How scenarios, ontology, actions, and the resolver fit together")
    _subtitle(slide,
              "Scenarios are the recipes. Ontology + KnowledgeResolver are how "
              "the agent finds data. Actions are how it writes back — always "
              "behind a human approval.")

    # Mini-diagram on the left (~58% width)
    slide.shapes.add_picture(
        str(COMPONENT_PNG),
        Inches(0.3), Inches(1.5),
        width=Inches(7.7),
    )

    # Right side: 4 component captions + fallback panel
    # ----- Top right: 4 caption cards in a 2x2 grid -----
    captions = [
        {
            "title": "Scenario",
            "fill": CYAN_FILL,
            "body": "A YAML recipe naming the stages, what to bind at each, "
                    "who reviews it, and what to execute on approve. Each "
                    "scenario references ontology classes — no raw SQL.",
            "tie": "SC-TC-007 binds Product, SanctionedEntity, PriorOverride.",
        },
        {
            "title": "Ontology + mappings",
            "fill": CYAN_FILL,
            "body": "Business classes (Supplier, Product, …) mapped to "
                    "physical columns in any registered source. Authored "
                    "once; reused by every scenario.",
            "tie": "SanctionedEntity → sanctions_csv.entity_name.",
        },
        {
            "title": "KnowledgeResolver",
            "fill": NAVY_FILL,
            "body": "The Protocol every connector implements. Ontology passes "
                    "it (class, filters) and gets back rows — without the "
                    "scenario ever knowing CSV vs Postgres vs Neo4j.",
            "tie": "Same Protocol serves CSV and the Neo4j supply graph.",
        },
        {
            "title": "Action + executor",
            "fill": AMBER_FILL,
            "body": "A typed write op with a named hand-written executor "
                    "(sql_update / http_request). The LLM only fills the "
                    "args; the executor only runs after reviewer approval.",
            "tie": "On approve: update_outcome_decided_by runs.",
        },
    ]

    # 2x2 grid on the right
    grid_left = Inches(8.15)
    grid_top = Inches(1.5)
    cell_w = Inches(2.5)
    cell_h = Inches(2.1)
    gap = Inches(0.1)
    for i, c in enumerate(captions):
        col = i % 2
        row = i // 2
        left = grid_left + (cell_w + gap) * col
        top = grid_top + (cell_h + gap) * row
        tf = _add_text_box(slide, left, top, cell_w, cell_h,
                           fill=c["fill"], border=True)
        p = tf.paragraphs[0]
        _set_run(p.add_run(), c["title"], size=12, bold=True, color=NAVY)
        p.space_after = Pt(3)
        p = tf.add_paragraph()
        _set_run(p.add_run(), c["body"], size=9.5, color=INK)
        p.space_after = Pt(4)
        p = tf.add_paragraph()
        _set_run(p.add_run(), c["tie"], size=8.5, italic=True, color=GREY,
                 name="Consolas")

    # ----- Bottom right: fallback chain panel -----
    fb_left = Inches(8.15)
    fb_top = Inches(5.85)
    fb_w = Inches(5.1)
    fb_h = Inches(1.5)
    tf = _add_text_box(slide, fb_left, fb_top, fb_w, fb_h,
                       fill=LIGHT_GREY, border=True)
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "Three composer paths  (one prompt, three lanes)",
             size=11, bold=True, color=NAVY)
    p.space_after = Pt(3)

    _bullet(tf,
            "Scenario match — Scenario + Ontology + Resolver  (chip path)",
            size=10, bold_lead=True)
    _bullet(tf,
            "Ontology lookup — Ontology + Resolver only  (no scenario; opt-in fallback)",
            size=10, bold_lead=True)
    _bullet(tf,
            "NL write action — Action + executor + HITL  (no scenario; opt-in fallback)",
            size=10, bold_lead=True)

    # Footer (left side, under the diagram)
    tf = _add_text_box(slide, Inches(0.3), Inches(7.0), Inches(7.7), Inches(0.45))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(),
             "Adding a new data source = registration + a mapping. Every "
             "existing scenario picks it up — no rewrites.",
             size=10, italic=True, color=GREY)


# =============================================================================
# Slide 4 — Walk-through of one request
# =============================================================================
def slide_walkthrough(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _title(slide, "How a request flows through the system")
    _subtitle(slide,
              "One real operator prompt, six steps. Each step names the layer "
              "from the diagram that activates and what it does.")

    # Top: the prompt as a quote bar
    tf = _add_text_box(slide, Inches(0.4), Inches(1.5), Inches(12.6),
                       Inches(0.65), fill=NAVY_FILL, border=True)
    p = tf.paragraphs[0]
    _set_run(p.add_run(),
             "Operator types:  ",
             size=11, bold=True, color=NAVY)
    r = p.add_run()
    _set_run(r,
             "“Override the SC-TC-001 sanctions block on order ORD-44216 "
             "— customer says they have an OFAC license.”",
             size=11, italic=True, color=NAVY)

    # Six steps as 2 rows x 3 cols
    steps = [
        {
            "n": "1",
            "layer": "Frontend → Orchestrator",
            "title": "Operator types it",
            "body": "The chat composer sends the prompt to the backend over REST. "
                    "An SSE channel opens so the operator can watch each stage's "
                    "facts populate live.",
            "fill": LIGHT_GREY,
        },
        {
            "n": "2",
            "layer": "Agent runtime  (LLM classifier)",
            "title": "Agent classifies the request",
            "body": "The LLM picks SC-TC-007 (Trade override — sanctioned "
                    "counterparty) from the scenario catalog and rephrases it. "
                    "Operator sees the rephrasing and clicks Confirm.",
            "fill": AMBER_FILL,
        },
        {
            "n": "3",
            "layer": "Scenarios + Ontology + Resolvers",
            "title": "Framework binds knowledge",
            "body": "Two stages run: intake pulls the active TC-override policy "
                    "and the agent's IAM scope. Proposal pulls product master, "
                    "the OFAC sanctions hit, and the contract value — each "
                    "fact tagged with which ontology class asked and which "
                    "source answered.",
            "fill": CYAN_FILL,
        },
        {
            "n": "4",
            "layer": "Orchestrator  (policy gate)",
            "title": "Policy says: needs a human",
            "body": "Sanctions overrides require named compliance officer "
                    "approval, so the case enters the review stage and the "
                    "orchestrator pauses. A Teams Adaptive Card is rendered with "
                    "all bound facts, prior similar overrides, and three SOP "
                    "excerpts retrieved by semantic search.",
            "fill": NAVY_FILL,
        },
        {
            "n": "5",
            "layer": "Reviewer surface  →  Action registry",
            "title": "Reviewer decides",
            "body": "The compliance officer reads the card and clicks Approve / "
                    "Reject / Need more info. On approve, the orchestrator wakes "
                    "the paused case and dispatches the write action to its "
                    "named executor (sql_update or http_request).",
            "fill": EMERALD_FILL,
        },
        {
            "n": "6",
            "layer": "Lineage recorder",
            "title": "Audit trail finalizes",
            "body": "Every fact bound, every query issued, every decision made: "
                    "appended to an immutable per-case log. Replayable later for "
                    "compliance review — even with a forced reviewer "
                    "decision to compare counterfactuals.",
            "fill": LIGHT_GREY,
        },
    ]

    # Layout: 3 columns x 2 rows
    cols = 3
    col_w = Inches(4.15)
    col_h = Inches(2.4)
    col_gap = Inches(0.1)
    row_gap = Inches(0.15)
    left_origin = Inches(0.4)
    top_origin = Inches(2.3)

    for i, step in enumerate(steps):
        col = i % cols
        row = i // cols
        left = Inches(0.4 + col * (4.15 + 0.1))
        top = Inches(2.3 + row * (2.4 + 0.15))
        tf = _add_text_box(slide, left, top, col_w, col_h,
                           fill=step["fill"], border=True)

        # Header line: number + layer tag
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, f"{step['n']}.  ", size=14, bold=True, color=NAVY)
        r = p.add_run()
        _set_run(r, step["layer"], size=10.5, bold=True, color=GREY)
        p.space_after = Pt(2)

        # Step title
        p = tf.add_paragraph()
        _set_run(p.add_run(), step["title"], size=13, bold=True, color=NAVY)
        p.space_after = Pt(4)

        # Body
        p = tf.add_paragraph()
        _set_run(p.add_run(), step["body"], size=10.5, color=INK)

    # Footer
    tf = _add_text_box(slide, Inches(0.4), Inches(7.15), Inches(12.6), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(),
             "Autonomous cases finish in ~10 seconds; HITL cases pause "
             "indefinitely at step 4 until a named reviewer decides.",
             size=10.5, italic=True, color=GREY)


# =============================================================================
# Entry
# =============================================================================
def build_deck(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_audiences(prs)
    slide_architecture(prs)
    slide_components(prs)
    slide_walkthrough(prs)

    prs.save(str(out))


if __name__ == "__main__":
    print(f"Rendering component diagram → {COMPONENT_PNG}")
    render_component_diagram(COMPONENT_PNG)
    print(f"Building deck → {OUT_PPTX}")
    build_deck(OUT_PPTX)
    print(f"Done. {OUT_PPTX.stat().st_size // 1024} KB")
