#!/usr/bin/env python3
"""Generate docs/datasets-and-rules.pptx — domain-expert reference deck.

Audience: procurement / logistics SMEs. Business-first language, no
JSON / YAML / SQL. Light tables only where structure helps.

Re-runnable: regenerates the file in-place. Scenario content is read
from backend/scenarios/*.yaml at build time so the deck doesn't drift
from the live catalogue.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = REPO_ROOT / "docs" / "datasets-and-rules.pptx"
SCENARIOS_DIR = REPO_ROOT / "backend" / "scenarios"


# Palette — matches the platform's indigo accent.
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x31, 0x27, 0xA8)
INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
GREY_DARK = RGBColor(0x1F, 0x29, 0x37)
GREY_BODY = RGBColor(0x37, 0x41, 0x51)
GREY_META = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# -----------------------------------------------------------------------------
# Slide helpers
# -----------------------------------------------------------------------------
def _new_blank_slide(prs: Presentation):
    """Blank slide layout — we paint everything ourselves."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _accent_bar(slide, *, color=INDIGO, top=Inches(0)) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _set_paragraph(p, text: str, *, size=14, bold=False, color=GREY_BODY,
                   align=PP_ALIGN.LEFT) -> None:
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def _textbox(slide, left, top, width, height) -> object:
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _new_blank_slide(prs)
    # Indigo background block on the left third
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = INDIGO
    block.line.fill.background()

    label = _textbox(slide, Inches(0.6), Inches(0.7), Inches(3.5), Inches(0.4))
    _set_paragraph(label.paragraphs[0], "KNOWLEDGE FABRIC",
                   size=11, bold=True, color=WHITE)

    title_tf = _textbox(slide, Inches(5.0), Inches(2.4), Inches(7.8), Inches(2.4))
    _set_paragraph(title_tf.paragraphs[0], title, size=40, bold=True,
                   color=GREY_DARK)
    p = title_tf.add_paragraph()
    _set_paragraph(p, subtitle, size=20, color=INDIGO_DARK)

    meta = _textbox(slide, Inches(5.0), Inches(5.5), Inches(7.8), Inches(1.5))
    _set_paragraph(meta.paragraphs[0],
                   "Reference deck for domain experts", size=14, color=GREY_META)
    p = meta.add_paragraph()
    _set_paragraph(p, "Datasets that feed the agent + rules the system enforces.",
                   size=14, color=GREY_META)


def add_section_divider(prs: Presentation, *, section: str, name: str,
                        subtitle: str) -> None:
    slide = _new_blank_slide(prs)
    # Full indigo background.
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = INDIGO
    bg.line.fill.background()

    eyebrow_tf = _textbox(slide, Inches(1.0), Inches(2.7), Inches(11), Inches(0.5))
    _set_paragraph(eyebrow_tf.paragraphs[0],
                   f"SECTION {section}", size=14, bold=True, color=WHITE)

    title_tf = _textbox(slide, Inches(1.0), Inches(3.2), Inches(11), Inches(1.4))
    _set_paragraph(title_tf.paragraphs[0], name, size=48, bold=True, color=WHITE)

    sub_tf = _textbox(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(1.5))
    _set_paragraph(sub_tf.paragraphs[0], subtitle, size=18, color=INDIGO_LIGHT)


def add_content_slide(
    prs: Presentation,
    *,
    eyebrow: str,
    title: str,
    body: list[str],
    side: list[tuple[str, str]] | None = None,
) -> None:
    """Standard content slide. Eyebrow + title + bullet body, with an
    optional side-panel of (label, value) pairs (e.g. "Used by: SC-PP-008")."""
    slide = _new_blank_slide(prs)
    _accent_bar(slide)

    eb = _textbox(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4))
    _set_paragraph(eb.paragraphs[0], eyebrow.upper(),
                   size=11, bold=True, color=INDIGO)

    tt = _textbox(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9))
    _set_paragraph(tt.paragraphs[0], title, size=30, bold=True, color=GREY_DARK)

    # Body column.
    body_width = Inches(7.6) if side else Inches(12.1)
    body_tf = _textbox(slide, Inches(0.6), Inches(2.1), body_width, Inches(5.0))
    for i, bullet in enumerate(body):
        para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        # Hand-rolled bullet — slide-layout bullets in pptx are clunky.
        _set_paragraph(para, "•  " + bullet, size=15, color=GREY_BODY)
        para.space_after = Pt(8)

    # Side panel — pairs.
    if side:
        panel_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.0),
            Inches(4.2), Inches(5.0),
        )
        panel_bg.fill.solid()
        panel_bg.fill.fore_color.rgb = INDIGO_LIGHT
        panel_bg.line.fill.background()

        panel_tf = _textbox(slide, Inches(8.8), Inches(2.2), Inches(3.8), Inches(4.6))
        for i, (label, value) in enumerate(side):
            label_p = panel_tf.paragraphs[0] if i == 0 else panel_tf.add_paragraph()
            _set_paragraph(label_p, label.upper(), size=10, bold=True, color=INDIGO_DARK)
            label_p.space_after = Pt(2)
            value_p = panel_tf.add_paragraph()
            _set_paragraph(value_p, value, size=12, color=GREY_DARK)
            value_p.space_after = Pt(10)


# -----------------------------------------------------------------------------
# Scenario loading
# -----------------------------------------------------------------------------
@dataclass
class ScenarioMeta:
    sid: str
    title: str
    domain: str
    autonomous: bool
    teams_headline: str
    interpreted_as: str
    reviewer: str
    keywords: list[str]
    raw: dict


def _load_scenarios() -> dict[str, ScenarioMeta]:
    out: dict[str, ScenarioMeta] = {}
    for path in sorted(SCENARIOS_DIR.glob("*.yaml")):
        data = yaml.safe_load(path.read_text(encoding="utf-8"))
        if not isinstance(data, dict):
            continue
        rev = data.get("reviewer_role") or {}
        out[data["id"]] = ScenarioMeta(
            sid=data["id"],
            title=data.get("title") or data["id"],
            domain=data.get("domain") or "",
            autonomous=bool(data.get("autonomous")),
            teams_headline=data.get("teams_headline") or "",
            interpreted_as=data.get("interpreted_as") or "",
            reviewer=f"{rev.get('label', '')} · {rev.get('name', '')}".strip(" ·"),
            keywords=list(data.get("match_keywords") or []),
            raw=data,
        )
    return out


# -----------------------------------------------------------------------------
# Deck content
# -----------------------------------------------------------------------------
def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    scenarios = _load_scenarios()

    # =========================================================================
    # Title + intro
    # =========================================================================
    add_title_slide(
        prs,
        title="Datasets & rules",
        subtitle="What the agent sees, and how decisions get made",
    )

    add_content_slide(
        prs,
        eyebrow="How to read this deck",
        title="Where to look, what to verify",
        body=[
            "The agent never invents data. Every fact on a case card "
            "comes from one of the data sources in Section A.",
            "Concepts the agent knows about (Supplier, Carrier, "
            "SanctionsProximity…) are defined in Section B's ontology.",
            "Scenarios in Section C are the work the agent can do. Each "
            "is triggered by an operator prompt and follows a fixed "
            "evidence-gathering recipe.",
            "Rules in Section D are the guardrails: when the agent acts "
            "alone, when the reviewer is asked, when the platform learns.",
            "If you spot a dataset that's wrong, a scenario that's missing "
            "a step, or a rule that doesn't match your domain — that's "
            "exactly the feedback we want.",
        ],
        side=[
            ("Audience", "Procurement & logistics domain experts"),
            ("Purpose", "Verify the agent's worldview against yours"),
            ("Out of scope", "Code, schemas, SQL — kept business-first"),
        ],
    )

    # =========================================================================
    # A. Data sources
    # =========================================================================
    add_section_divider(
        prs, section="A", name="Data sources",
        subtitle=(
            "Everything the agent shows you started as a row in one of these. "
            "If something is missing or stale here, it'll be missing in the "
            "case too."
        ),
    )

    add_content_slide(
        prs,
        eyebrow="Section A · overview",
        title="14 data sources, four kinds",
        body=[
            "Spreadsheets (CSV) for hand-curated master data: products, "
            "carriers, lanes, ports, rate cards.",
            "Operational stores (SQLite) for live state that changes when "
            "the agent acts: bookings, demurrage charges, prior decisions.",
            "Vector store for free-text policy documents — semantic search "
            "over SOP excerpts.",
            "Graph (Neo4j) for relationships: ownership chains, carrier "
            "alliances, sanctions proximity walks.",
            "A single test API (jsonplaceholder.com) confirms the HTTP "
            "connector works without needing real keys.",
        ],
        side=[
            ("Where it lives", "backend/data/sources.yaml"),
            ("Who owns curation", "Domain SMEs (you)"),
            ("Refresh cadence", "On-demand for now"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="A · spreadsheets",
        title="Procurement master data",
        body=[
            "Products master (15 SKUs incl. encryption-classified items). "
            "Used by supplier onboarding to check what the candidate "
            "actually makes.",
            "Sanctions list (~30 OFAC-style entities). Schema matches the "
            "real SDN feed so swapping in the live one is a config change.",
            "Both are CSVs an SME can edit in Excel. Changes take effect "
            "next time the backend reloads.",
        ],
        side=[
            ("Files", "products.csv\nsanctions_sdn_sample.csv"),
            ("Backs ontology", "tcs_core.Product\ntcs_core.SanctionedEntity"),
            ("Demo scenarios", "SC-PP-008, SC-TC-008"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="A · spreadsheets",
        title="Logistics master data",
        body=[
            "Carriers (25 rows, real SCAC codes) — who can move freight, "
            "their alliance, their reliability score.",
            "Lanes (40 origin-destination corridors) — the routes shipments "
            "travel on.",
            "Ports (40 UN/LOCODEs with live operational state) — "
            "congestion, closures.",
            "Freight rates (~60 rows: lane × carrier × container type, "
            "mix of contract + spot pricing).",
            "Customer orders (~750 rows) — drives the SLA-risk scenarios "
            "via the penalty_per_day_usd field.",
        ],
        side=[
            ("Files", "carriers.csv, lanes.csv,\nports.csv, freight_rates.csv,\ncustomer_orders.csv"),
            ("Refresh", "Hand-curated; future: pull from a live OMS"),
            ("Demo scenarios", "SC-LN-002 mode switch,\nSC-LN-PORT-DISRUPTION-020"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="A · spreadsheets",
        title="Live shipment activity",
        body=[
            "Shipments (~3,000 rows) — every shipment we've booked + its "
            "live status snapshot.",
            "Shipment events (~19k rows, ~6 per shipment) — milestone + "
            "exception stream (gate-out, vessel-load, arrived, "
            "demurrage-accrual, etc.).",
            "Per-SKU performance aggregates (3 quarters × 15 SKUs) — used "
            "by the carrier scorecard scenario.",
            "These mirror what a live tracking system feed would look "
            "like; the agent reads them in the same shape.",
        ],
        side=[
            ("Files", "shipments.csv,\nshipment_events.csv,\nshipment_performance.csv"),
            ("Volume", "~22k rows total"),
            ("Demo scenarios", "Status lookup, trace,\nport disruption"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="A · operational stores",
        title="Write targets: where the agent acts",
        body=[
            "Logistics SQLite holds bookings (the active state of every "
            "shipment) and demurrage charges. When the agent re-tenders "
            "a booking or waives a charge, it writes here.",
            "Governance SQLite holds prior reviewer decisions. The agent "
            "reads this to show \"here's how reviewers handled similar "
            "cases before.\"",
            "Both are append-mostly — old rows are preserved so the "
            "audit trail stays honest.",
            "A Postgres mirror of governance exists for swap-in if SQLite "
            "ever becomes a bottleneck. Same schema.",
        ],
        side=[
            ("Files", "logistics.sqlite,\ngovernance.sqlite"),
            ("Write actions", "retender_booking,\nwaive_demurrage"),
            ("Backs ontology", "logistics.Booking,\ntcs_core.PriorOverride"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="A · semantic + graph",
        title="Policy library and the supply-chain graph",
        body=[
            "Policy corpus: 5 markdown SOP excerpts indexed for semantic "
            "search. When a scenario asks \"is this onboarding consistent "
            "with our SOP?\", the agent pulls the most-relevant excerpts.",
            "Neo4j graph: 280 nodes covering supplier corporate networks, "
            "carrier alliances, and ownership chains. Lets the agent "
            "answer \"is this supplier 3 hops from a sanctioned entity?\" "
            "with one query.",
            "The graph is where the heavy lifting happens for "
            "SC-PP-008 (sanctions proximity), SC-PP-UBO-023 (ownership "
            "walk), SC-PP-CARRIER-EXP-022 (alliance reach).",
            "Neo4j is also why \"5 sources fused in one case\" works — "
            "graph queries traverse + return enriched rows in one shot.",
        ],
        side=[
            ("Files", "policy_corpus/ (markdown),\nNeo4j (external service)"),
            ("Backs ontology", "tcs_core.PolicyExcerpt,\nsupply_chain.* (most)"),
            ("Demo scenarios", "SC-PP-008,\nSC-PP-UBO-023,\nSC-PP-CARRIER-EXP-022"),
        ],
    )

    # =========================================================================
    # B. Ontology
    # =========================================================================
    add_section_divider(
        prs, section="B", name="Ontology",
        subtitle=(
            "The vocabulary the agent uses. Concepts the platform "
            "understands — Supplier, Carrier, SanctionsProximity — and "
            "where each concept's data lives."
        ),
    )

    add_content_slide(
        prs,
        eyebrow="Section B · overview",
        title="Why ontology matters to a domain expert",
        body=[
            "Every fact card in a case has a class label (e.g. "
            "SanctionsProximity). The ontology is the dictionary that "
            "says what each class means.",
            "Ontology lives separately from data sources so the same "
            "concept (Supplier) can come from CSV today and the live ERP "
            "tomorrow without scenarios needing to change.",
            "If a concept is missing — \"we also care about beneficial "
            "owners' citizenship\" — that's an ontology change, not a "
            "data change.",
            "Four families today: tcs_core (cross-domain), supply_chain "
            "(procurement-side), logistics (movement-side), and a "
            "demo-only variant for sandbox work.",
        ],
        side=[
            ("Total classes", "27 across 4 families"),
            ("Where it lives", "backend/ontologies/*.yaml"),
            ("Editable by", "Engineering + SME pairing"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="B · tcs_core",
        title="Cross-domain core — used by every scenario",
        body=[
            "Product — anything we buy. Hooks into the products spreadsheet.",
            "SanctionedEntity — a single hit from the OFAC list. Used "
            "for direct sanctions checks.",
            "PriorOverride — a past reviewer decision. The agent uses "
            "these as precedent (\"reviewers rejected 4 similar cases\").",
            "PolicyExcerpt — a chunk of SOP text retrievable via "
            "semantic search.",
        ],
        side=[
            ("Used by", "All HITL scenarios"),
            ("Owned by", "Cross-functional"),
            ("4 classes", "Product, SanctionedEntity,\nPriorOverride, PolicyExcerpt"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="B · supply_chain",
        title="Procurement vocabulary",
        body=[
            "Supplier — the candidate or active vendor. Backed by the graph.",
            "PurchaseOrder — what we've bought from them.",
            "SanctionsProximity — a walk from a supplier through ownership "
            "edges to any sanctioned entity. Hop distance matters.",
            "UltimateBeneficialOwner — the human(s) ultimately behind a "
            "corporate chain. Used for tier-1 onboarding scrutiny.",
            "CarrierExposure — when an SDN listing affects a carrier we "
            "use indirectly via supplier choice.",
            "AlternativeSupplier and TierSupplier — for resilience scenarios "
            "(primary blocked, multi-tier visibility).",
        ],
        side=[
            ("Where its data lives", "Mostly Neo4j graph,\nplus suppliers.csv"),
            ("Used by", "All SC-PP-* scenarios"),
            ("10 classes", "(listed in body)"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="B · logistics",
        title="Movement vocabulary",
        body=[
            "Shipment + ShipmentEvent — what's moving and what's happened "
            "to it.",
            "Carrier + Lane + Port — the transport network.",
            "Booking + BookingHistory — our active commitments (the write "
            "target for re-tender actions).",
            "FreightRate — what we'd pay today on a given route, by carrier "
            "and container type.",
            "CustomerOrder — the demand we're fulfilling, including SLA + "
            "penalty terms.",
            "DemurrageCharge — penalty billed by ports when containers sit "
            "too long. Targeted by the auto-dispute scenario.",
        ],
        side=[
            ("Where its data lives", "CSVs + logistics.sqlite"),
            ("Used by", "All SC-LN-* scenarios"),
            ("11 classes", "(listed in body)"),
        ],
    )

    # =========================================================================
    # C. Scenarios
    # =========================================================================
    add_section_divider(
        prs, section="C", name="Scenarios",
        subtitle=(
            "Each scenario is a recipe the agent follows: triggered by a "
            "prompt, gathers evidence from the data sources via the "
            "ontology, and presents a decision to either a reviewer (HITL) "
            "or applies a guardrail (autonomous)."
        ),
    )

    add_content_slide(
        prs,
        eyebrow="Section C · overview",
        title="Anatomy of a scenario",
        body=[
            "Triggered by match_keywords — phrases the operator might "
            "type. The agent picks the highest-confidence scenario.",
            "Runs three stages: (1) Policy & scope — what the agent is "
            "supposed to do, (2) Agent gathered — facts pulled from data "
            "sources, (3) For your decision — what the reviewer sees.",
            "Each stage's facts come from ontology_queries: the platform "
            "version of \"go look up X from source Y where Z\".",
            "Settles via rationale_reasons (canned options the reviewer "
            "picks) + closing_messages (what the operator and audit log "
            "see).",
        ],
        side=[
            ("Total shipped", "28 scenarios"),
            ("By domain", "7 procurement,\n11 logistics,\n2 trade compliance,\n8 ontology lookups"),
            ("HITL vs autonomous", "21 HITL, 7 autonomous"),
        ],
    )

    # ---- Procurement domain ------------------------------------------------
    add_section_divider(
        prs, section="C-PP", name="Procurement scenarios",
        subtitle=(
            "7 scenarios across supplier onboarding, ownership disclosure, "
            "alternative-supplier discovery, and trade-compliance "
            "containment."
        ),
    )

    for sid in ("SC-PP-008", "SC-PP-UBO-023", "SC-PP-CARRIER-EXP-022"):
        if sid in scenarios:
            _add_scenario_spotlight(prs, scenarios[sid])

    add_content_slide(
        prs,
        eyebrow="C-PP · the rest",
        title="Other procurement scenarios",
        body=[
            "SC-PP-007 — Supplier onboarding (legacy single-source variant; "
            "kept for comparison vs SC-PP-008's multi-source review).",
            "SC-PP-ALT-SUP-024 — Alternative-supplier discovery when the "
            "primary supplier is blocked (sanctions, capacity, dispute).",
            "SC-PP-AUTO-014 — Reorder-point adjustment within the "
            "autonomous envelope (small SKU stock changes the agent makes "
            "without HITL).",
            "SC-PP-TIER-N-025 — Multi-tier supply-chain visibility — "
            "agent walks tier-2 and tier-3 suppliers for resilience risk.",
        ],
        side=[
            ("Type breakdown", "6 HITL, 1 autonomous"),
            ("Who reviews", "Procurement lead / category manager"),
            ("Demo-ready", "SC-PP-008 is the canonical one"),
        ],
    )

    # ---- Logistics domain --------------------------------------------------
    add_section_divider(
        prs, section="C-LN", name="Logistics scenarios",
        subtitle=(
            "11 scenarios covering mode switches, port disruptions, "
            "carrier scorecards, demurrage disputes, shipment lookups + "
            "traces, and SLA-risk routing decisions."
        ),
    )

    for sid in ("SC-LN-002", "SC-LN-PORT-DISRUPTION-020", "SC-LN-DEMURRAGE-AUTO-025"):
        if sid in scenarios:
            _add_scenario_spotlight(prs, scenarios[sid])

    add_content_slide(
        prs,
        eyebrow="C-LN · the rest",
        title="Other logistics scenarios",
        body=[
            "SC-LN-CARRIER-SCORECARD-019 — One-stop carrier overview "
            "(scorecard + recent shipments). Autonomous lookup.",
            "SC-LN-DEMURRAGE-DISPUTE-024 — Manual demurrage dispute when "
            "the charge exceeds the autonomous envelope.",
            "SC-LN-ETA-PUSH-023 — Push updated ETA to the customer system "
            "(autonomous).",
            "SC-LN-SLA-RISK-022 — Customer SLA risk: expedite at extra "
            "cost or accept the penalty.",
            "SC-LN-STATUS-009 / SC-LN-TRACE-018 — Shipment status + full "
            "event-history lookups (autonomous).",
            "SC-LN-TENDER-021 — Re-tender a booking to a new carrier "
            "(HITL).",
            "SC-LN-DEMO-RESET-999 — Demo-tooling only: resets booking "
            "state between demos.",
        ],
        side=[
            ("Type breakdown", "5 HITL, 6 autonomous"),
            ("Who reviews HITL", "Logistics planner /\nplanner lead"),
            ("Demo-ready", "SC-LN-002, port disruption,\ndemurrage auto"),
        ],
    )

    # ---- Trade compliance + ontology lookups -------------------------------
    add_content_slide(
        prs,
        eyebrow="C-TC · trade compliance",
        title="Trade-compliance scenarios",
        body=[
            "SC-TC-007 — Trade override for a sanctioned counterparty. "
            "Compliance officer reviews the agent's evidence and decides "
            "whether to allow or block the transaction.",
            "SC-TC-008 — Live-data variant focused on sanctions hits "
            "specifically on encryption-classified products.",
            "Both feed off products_csv + sanctions_csv + the policy "
            "corpus. Decision is recorded against the customer-facing "
            "audit trail.",
        ],
        side=[
            ("Type", "Both HITL"),
            ("Reviewer", "Trade compliance officer"),
            ("Demo-ready", "SC-TC-008 with the encryption SKU"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="C-ONTO · self-service lookups",
        title="Ontology-lookup scenarios (auto-generated)",
        body=[
            "Every class with a binding gets a SC-ONTO-{ontology}-{class} "
            "lookup scenario, generated from the ontology itself.",
            "Lets an operator type \"look up products\" and the agent "
            "returns rows from products_csv via the tcs_core.Product class.",
            "Autonomous (no reviewer); the goal is fast self-service "
            "data exploration, not a decision.",
            "Examples: SC-ONTO-tcs_core-Product, SC-ONTO-supply_chain-Supplier, "
            "SC-ONTO-tcs_core-SanctionedEntity.",
        ],
        side=[
            ("Count", "8 auto-generated"),
            ("Type", "All autonomous"),
            ("Owned by", "Engineering — admins can add by clicking 'generate' on a class"),
        ],
    )

    # =========================================================================
    # D. Rules
    # =========================================================================
    add_section_divider(
        prs, section="D", name="Rules",
        subtitle=(
            "The guardrails: when the agent acts alone, when a human is "
            "asked, when the platform learns from past behaviour, and "
            "what stays immutable for the audit trail."
        ),
    )

    add_content_slide(
        prs,
        eyebrow="D · autonomous envelopes",
        title="When the agent acts alone",
        body=[
            "Each scenario carries an autonomous flag. True = the agent "
            "executes the action without HITL once it's gathered the "
            "evidence.",
            "Autonomous scenarios still have a guardrail: a policy "
            "(e.g. \"only for demurrage charges under $500 with a "
            "prior-win precedent\") that must pass.",
            "If the guardrail fails, the autonomous scenario routes to "
            "HITL automatically. Operator sees a guardrail-fired notice "
            "card.",
            "All 7 autonomous scenarios today are read-only lookups or "
            "very low-stakes writes (ETA push, demurrage auto-dispute).",
            "Adding a new autonomous scenario should be a sign-off "
            "conversation, not a quiet engineering change.",
        ],
        side=[
            ("Autonomous today", "7 of 28 scenarios"),
            ("Highest-stakes", "Demurrage auto-dispute\n(envelope: <$500 + prior wins)"),
            ("Default for new", "HITL until proven safe"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="D · promoted patterns (Phase 3b)",
        title="When the platform learns from reviewers",
        body=[
            "Reviewers can flag up to 3 facts as load-bearing at "
            "decision time. The platform stores these as structured signal.",
            "When the same driver (e.g. SanctionsProximity) shows up in "
            "many reject decisions on the same scenario, an admin can "
            "promote the pattern into the scenario itself.",
            "Promoted patterns surface as an advisory chip on future "
            "cases: \"in 8 of 8 similar cases, reviewers rejected for "
            "this driver.\" The reviewer is not bound by the suggestion.",
            "Guardrails: minimum 2 cases + 50% share of decision kind by "
            "default. Re-checked at promote time, not from the UI's "
            "cached view.",
            "Every promote and demote is a new scenario version. Audit "
            "trail records who promoted, when, from how many cases.",
        ],
        side=[
            ("Where it surfaces", "Insights modal (admin only),\nadvisory chip on cases"),
            ("Reversible", "Yes — demote button"),
            ("Audit-traced", "Yes — every change is a\nnew scenario version"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="D · safety rails",
        title="What stays immutable",
        body=[
            "Every case records the exact scenario version it ran "
            "against. Editing a scenario later doesn't retroactively "
            "change historical cases.",
            "The audit lineage is append-only. Stage events are written "
            "once and never edited or deleted.",
            "Scenario edits create a new version snapshot before the "
            "live file is overwritten. Older versions are inspectable in "
            "the Version history modal.",
            "Reviewer decisions are immutable once submitted. The only "
            "way to change a decision is to start a new replay case "
            "(which gets its own audit trail).",
            "Operator-authored YAML comments survive migration "
            "byte-for-byte — policy intent doesn't get silently erased.",
        ],
        side=[
            ("Why this matters", "Auditors can answer\n\"what spec was the agent\nrunning when this decision\nwas made?\""),
            ("Crash-safe", "Yes — mid-save crashes\nheal forward on reboot"),
            ("Concurrent-edit safe", "Yes — serialized writes"),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="D · access control",
        title="Who can do what",
        body=[
            "Operators submit cases and view their own. They cannot "
            "review their own cases or see other operators' cases.",
            "Reviewers (e.g. procurement lead, planner lead) approve / "
            "reject / request more info on cases routed to them. They "
            "see the full evidence package.",
            "Admins (a smaller group) edit scenarios, view the Insights "
            "modal, promote/demote patterns. Their actions all leave a "
            "trail.",
            "Auditors are admins with read-only intent — same surface, "
            "but the expectation is they don't promote or edit.",
        ],
        side=[
            ("Today", "Role enforced at the API\nboundary (403 on mismatch)"),
            ("Future", "Per-scenario role gating\n(Phase 4 candidate)"),
            ("Audit", "Every state change\ntied to a username"),
        ],
    )

    # =========================================================================
    # E. Glossary
    # =========================================================================
    add_section_divider(
        prs, section="E", name="Glossary",
        subtitle=(
            "Plain-English definitions for the words this deck uses. "
            "Useful as a back-of-the-room reference during a walk-through."
        ),
    )

    add_content_slide(
        prs,
        eyebrow="E · the basics",
        title="Words used in this deck",
        body=[
            "Case — one operator-initiated unit of work. Has a prompt, "
            "a scenario, evidence, and (eventually) a decision.",
            "Scenario — a named recipe for how the agent should handle a "
            "kind of work (\"onboard a supplier,\" \"switch a shipment mode\").",
            "Ontology — the dictionary of concepts (Supplier, Carrier, "
            "etc.) the platform understands.",
            "Data source — where rows of data live (a CSV, a SQLite "
            "table, the Neo4j graph, a vector store of SOP text).",
            "Lineage — the immutable record of every step the agent "
            "took during a case. The audit trail.",
            "Scenario version — every edit produces a new version; old "
            "versions stay forever for audit.",
        ],
        side=[
            ("If you forget one", "Open this slide.\nIt's the canonical list."),
        ],
    )

    add_content_slide(
        prs,
        eyebrow="E · Phase 2 + 3 words",
        title="Compounding-loop vocabulary",
        body=[
            "Load-bearing fact — one of the 1–3 evidence cards the "
            "reviewer flags as having tipped their decision.",
            "Highlight — the captured (source, ontology_type, id) "
            "reference for a load-bearing fact.",
            "Pattern — across many cases, a recurring "
            "(scenario, decision, ontology_type) combination flagged "
            "as load-bearing.",
            "Promote — admin action that writes a pattern into the "
            "scenario YAML so future cases see an advisory chip.",
            "Demote — the inverse of promote. Reversible audit trail.",
            "Advisory chip — the banner above the evidence map on a "
            "case where a promoted pattern fires.",
        ],
        side=[
            ("Where these live", "Section D · rules,\nand the Insights modal"),
            ("Why they matter", "This is the platform\nlearning over time."),
        ],
    )

    # Closing.
    add_section_divider(
        prs, section="•", name="Questions, gaps, corrections",
        subtitle=(
            "If a dataset is wrong, a scenario doesn't match how your "
            "team actually works, or a rule needs tightening — that's "
            "the highest-value feedback you can give us. Mark it on the "
            "deck, send it back, and we'll fold it in."
        ),
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


def _add_scenario_spotlight(prs: Presentation, m: ScenarioMeta) -> None:
    """Single-scenario detail slide."""
    decision_mode = "Autonomous" if m.autonomous else "Human-in-the-loop"
    # First N match keywords to keep the side panel compact.
    kw_preview = ", ".join(m.keywords[:5]) + (
        f" (+ {len(m.keywords) - 5} more)" if len(m.keywords) > 5 else ""
    )
    add_content_slide(
        prs,
        eyebrow=f"C · spotlight · {m.sid}",
        title=m.title,
        body=[
            f"What the agent does: {m.interpreted_as[:240]}",
            f"Decision mode: {decision_mode}. "
            + (f"Reviewer: {m.reviewer}." if not m.autonomous and m.reviewer else
               "No human review — autonomous via guardrail."),
            f"Teams headline (notification): \"{m.teams_headline}\"",
            f"Triggered by phrases like: {kw_preview}",
        ],
        side=[
            ("Scenario id", m.sid),
            ("Domain", m.domain),
            ("Source file", f"backend/scenarios/{m.sid}.yaml"),
        ],
    )


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
